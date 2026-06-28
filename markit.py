import sys
import re
import shutil
from pathlib import Path
from markitdown import MarkItDown
from llm_docx_markitdown import LlmDocxConverter
from llm_pdf_markitdown import LlmPdfConverter
from llm_pptx_markitdown import LlmPptxConverter


def extract_images_from_pptx(file_path, output_dir):
    """Extract images from PPTX file."""
    try:
        import pptx
    except ImportError:
        print("Warning: python-pptx not installed, skipping image extraction")
        return []

    images = []
    presentation = pptx.Presentation(file_path)
    image_counter = 1

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                try:
                    # Get image data
                    image = shape.image
                    image_bytes = image.blob

                    # Determine extension from content type or filename
                    ext = '.jpg'
                    if image.content_type:
                        if 'png' in image.content_type:
                            ext = '.png'
                        elif 'jpeg' in image.content_type or 'jpg' in image.content_type:
                            ext = '.jpg'
                        elif 'gif' in image.content_type:
                            ext = '.gif'
                    elif image.filename:
                        ext = Path(image.filename).suffix or '.jpg'

                    # Save image
                    image_filename = f"image_{image_counter:03d}{ext}"
                    image_path = output_dir / image_filename
                    image_path.write_bytes(image_bytes)

                    images.append(image_filename)
                    image_counter += 1
                except Exception as e:
                    print(f"Warning: Failed to extract image: {e}")

    return images


def extract_images_from_pdf(file_path, output_dir):
    """Extract images from PDF file."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        print("Warning: PyMuPDF not installed, skipping image extraction")
        return []

    images = []
    doc = fitz.open(file_path)
    image_counter = 1

    for page_num in range(len(doc)):
        page = doc[page_num]
        image_list = page.get_images()

        for img_index, img in enumerate(image_list):
            try:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]

                # Save image
                image_filename = f"image_{image_counter:03d}.{image_ext}"
                image_path = output_dir / image_filename
                image_path.write_bytes(image_bytes)

                images.append(image_filename)
                image_counter += 1
            except Exception as e:
                print(f"Warning: Failed to extract image: {e}")

    doc.close()
    return images


def extract_images_from_docx(file_path, output_dir):
    """Extract images from DOCX file."""
    try:
        import docx
    except ImportError:
        print("Warning: python-docx not installed, skipping image extraction")
        return []

    images = []
    try:
        doc = docx.Document(file_path)
    except Exception as e:
        print(f"Warning: Failed to read DOCX images: {e}")
        return []
    image_counter = 1

    # Extract from inline shapes
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                image_bytes = rel.target_part.blob

                # Determine extension from content type
                ext = '.jpg'
                if rel.target_part.content_type:
                    if 'png' in rel.target_part.content_type:
                        ext = '.png'
                    elif 'jpeg' in rel.target_part.content_type or 'jpg' in rel.target_part.content_type:
                        ext = '.jpg'
                    elif 'gif' in rel.target_part.content_type:
                        ext = '.gif'

                # Save image
                image_filename = f"image_{image_counter:03d}{ext}"
                image_path = output_dir / image_filename
                image_path.write_bytes(image_bytes)

                images.append(image_filename)
                image_counter += 1
            except Exception as e:
                print(f"Warning: Failed to extract image: {e}")

    return images


def update_markdown_image_paths(markdown_text, images, output_dir_relative):
    """Update image references in markdown to point to extracted images."""
    if not images:
        return re.sub(r'!\[([^\]]*)\]\(\)', "", markdown_text)

    # Find all image references in markdown
    image_pattern = r'!\[([^\]]*)\]\(([^\)]*)\)'
    matches = list(re.finditer(image_pattern, markdown_text))

    # Replace image paths with extracted image paths
    updated_markdown = markdown_text
    image_index = 0

    for match in matches:
        alt_text = match.group(1)
        old_path = match.group(2)

        # Skip if it's already a data URI or absolute path
        if old_path.startswith('data:') or old_path.startswith('http'):
            continue

        if image_index >= len(images):
            updated_markdown = updated_markdown.replace(match.group(0), "", 1)
            continue

        # Replace with new path
        new_path = f"{output_dir_relative}/{images[image_index]}"
        old_match = match.group(0)
        new_match = f"![{alt_text}]({new_path})"
        updated_markdown = updated_markdown.replace(old_match, new_match, 1)

        image_index += 1

    return updated_markdown


def process_document(file_path):
    """Process document: convert to markdown and extract images."""
    file_path = Path(file_path)

    if not file_path.exists():
        print(f"Error: File not found: {file_path}")
        sys.exit(1)

    # Determine output directory name
    file_stem = file_path.stem
    file_ext = file_path.suffix.lower().lstrip('.')
    output_dir_name = f"{file_stem}_{file_ext}"

    # Create output directory structure: results/{filename}_{ext}/
    result_dir = Path("./results") / output_dir_name
    images_dir = result_dir / "images"
    if images_dir.exists():
        shutil.rmtree(images_dir)
    images_dir.mkdir(parents=True, exist_ok=True)

    # Extract images based on file type
    images = []
    if file_ext == 'doc':
        images = extract_images_from_docx(file_path, images_dir)

    # Convert to markdown using MarkItDown
    md = MarkItDown(enable_plugins=False)
    md.register_converter(LlmDocxConverter(), priority=-1)
    md.register_converter(LlmPdfConverter(), priority=-1)
    md.register_converter(LlmPptxConverter(), priority=-1)
    convert_kwargs = {}
    if file_ext in ['docx', 'pdf', 'pptx']:
        convert_kwargs["image_output_dir"] = str(images_dir)
        convert_kwargs["image_output_prefix"] = "./images"
    result = md.convert(str(file_path), **convert_kwargs)

    if file_ext in ['docx', 'pdf', 'pptx']:
        updated_markdown = re.sub(r'!\[([^\]]*)\]\(\)', "", result.markdown)
        images = sorted(path.name for path in images_dir.iterdir() if path.is_file())
    else:
        # Update markdown with extracted image paths (relative to content.md)
        updated_markdown = update_markdown_image_paths(result.markdown, images, "./images")

    # Save markdown to content.md
    content_file = result_dir / "content.md"
    content_file.write_text(updated_markdown, encoding='utf-8')

    # Print summary
    print(f"✓ Converted: {file_path}")
    print(f"✓ Markdown saved to: {content_file}")
    print(f"✓ Extracted {len(images)} images to: {images_dir}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python markit.py <file_path>")
        print("Example: python markit.py ./files/ll.pptx")
        sys.exit(1)

    file_path = sys.argv[1]
    process_document(file_path)
